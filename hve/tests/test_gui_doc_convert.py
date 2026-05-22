"""test_gui_doc_convert.py — hve.gui.doc_convert の純粋ユニットテスト。

PySide6 / Qt に依存しないため、PySide6 未インストール環境でも実行可能。
任意依存（markitdown）が無くてもテキスト系（md/txt/csv）は動作する。
html/docx/pdf/xlsx/xls/pptx は markitdown インストール時のみ検証。
"""

from __future__ import annotations

import tempfile
import unittest
from pathlib import Path


def _has_markitdown() -> bool:
    """markitdown が import 可能かを返す（テストの skip 判定用）。"""
    try:
        import markitdown  # type: ignore[import-not-found]  # noqa: F401
        return True
    except ImportError:
        return False


class TestSupportedExtensions(unittest.TestCase):
    def test_includes_common_formats(self) -> None:
        from hve.gui.doc_convert import supported_extensions

        exts = set(supported_extensions())
        for required in (".md", ".markdown", ".txt", ".csv"):
            self.assertIn(required, exts)

    def test_includes_markitdown_formats(self) -> None:
        from hve.gui.doc_convert import supported_extensions

        exts = set(supported_extensions())
        for required in (".html", ".htm", ".docx", ".pdf", ".xlsx", ".xls", ".pptx"):
            self.assertIn(required, exts)

    def test_is_supported(self) -> None:
        from hve.gui.doc_convert import is_supported

        self.assertTrue(is_supported(Path("foo.md")))
        self.assertTrue(is_supported(Path("foo.txt")))
        self.assertTrue(is_supported(Path("FOO.CSV")))  # 大文字も OK
        self.assertFalse(is_supported(Path("foo.exe")))
        self.assertFalse(is_supported(Path("foo")))


class TestSafeFilename(unittest.TestCase):
    def test_ascii_pass_through(self) -> None:
        from hve.gui.doc_convert import safe_filename

        self.assertEqual(safe_filename("hello.txt"), "hello.md")

    def test_spaces_replaced(self) -> None:
        from hve.gui.doc_convert import safe_filename

        self.assertEqual(safe_filename("Hello World.txt"), "Hello_World.md")

    def test_special_chars_replaced(self) -> None:
        from hve.gui.doc_convert import safe_filename

        # `,` `(` `)` などは `_` に置換
        result = safe_filename("Business Plan (2026), v2.pdf")
        # 連続する `_` も処理される
        self.assertTrue(result.endswith(".md"))
        # スペース・カンマ・括弧は全て `_` に置換される
        for forbidden in (" ", ",", "(", ")"):
            self.assertNotIn(forbidden, result)

    def test_unicode_normalized(self) -> None:
        """日本語ファイル名は `_` に置換される（ASCII 安全化、設計書 §7.4）。"""
        from hve.gui.doc_convert import safe_filename

        # 純粋に日本語のみのファイル名は "file.md" にフォールバック
        result = safe_filename("メモ.txt")
        self.assertEqual(result, "file.md")


class TestConvertMd(unittest.TestCase):
    def test_md_passthrough(self) -> None:
        from hve.gui.doc_convert import convert_file

        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "src.md"
            src.write_text("# Hello\n\nWorld", encoding="utf-8")
            out_dir = Path(td) / "out"
            r = convert_file(src, out_dir=out_dir)
            self.assertTrue(r.ok)
            self.assertIsNotNone(r.converted_path)
            assert r.converted_path is not None
            content = r.converted_path.read_text(encoding="utf-8")
            self.assertIn("# Hello", content)
            self.assertIn("World", content)


class TestConvertTxt(unittest.TestCase):
    def test_txt_to_markdown(self) -> None:
        from hve.gui.doc_convert import convert_file

        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "notes.txt"
            src.write_text("This is plain text.\nLine 2.", encoding="utf-8")
            out_dir = Path(td) / "out"
            r = convert_file(src, out_dir=out_dir)
            self.assertTrue(r.ok, msg=f"error: {r.error}")
            assert r.converted_path is not None
            content = r.converted_path.read_text(encoding="utf-8")
            # 見出しが付与される
            self.assertIn("# notes", content)
            self.assertIn("This is plain text.", content)


class TestConvertCsv(unittest.TestCase):
    def test_csv_to_markdown_table(self) -> None:
        from hve.gui.doc_convert import convert_file

        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "data.csv"
            src.write_text("Name,Age\nAlice,30\nBob,25", encoding="utf-8")
            out_dir = Path(td) / "out"
            r = convert_file(src, out_dir=out_dir)
            self.assertTrue(r.ok)
            assert r.converted_path is not None
            content = r.converted_path.read_text(encoding="utf-8")
            self.assertIn("| Name | Age |", content)
            self.assertIn("| --- | --- |", content)
            self.assertIn("| Alice | 30 |", content)

    def test_csv_with_pipe_escapes(self) -> None:
        from hve.gui.doc_convert import convert_file

        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "pipe.csv"
            src.write_text("col1,col2\nval|with|pipe,normal", encoding="utf-8")
            out_dir = Path(td) / "out"
            r = convert_file(src, out_dir=out_dir)
            self.assertTrue(r.ok)
            assert r.converted_path is not None
            content = r.converted_path.read_text(encoding="utf-8")
            # `|` がエスケープされる
            self.assertIn(r"val\|with\|pipe", content)


class TestConvertUnsupported(unittest.TestCase):
    def test_unsupported_extension(self) -> None:
        from hve.gui.doc_convert import convert_file

        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "binary.exe"
            src.write_bytes(b"\x00\x01\x02")
            out_dir = Path(td) / "out"
            r = convert_file(src, out_dir=out_dir)
            self.assertFalse(r.ok)
            self.assertIsNotNone(r.error)
            assert r.error is not None
            self.assertIn("未対応", r.error)

    def test_nonexistent_file(self) -> None:
        from hve.gui.doc_convert import convert_file

        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "does-not-exist.txt"
            out_dir = Path(td) / "out"
            r = convert_file(src, out_dir=out_dir)
            self.assertFalse(r.ok)
            assert r.error is not None
            self.assertIn("存在しません", r.error)


@unittest.skipIf(_has_markitdown(), "markitdown is installed; missing-dep path not exercised")
class TestMarkItDownMissing(unittest.TestCase):
    """markitdown 未インストール環境で明示的なエラーメッセージが返る。"""

    def test_html_missing_markitdown_message(self) -> None:
        from hve.gui.doc_convert import convert_file

        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "page.html"
            src.write_text("<h1>Hello</h1>", encoding="utf-8")
            out_dir = Path(td) / "out"
            r = convert_file(src, out_dir=out_dir)
            self.assertFalse(r.ok)
            assert r.error is not None
            self.assertIn("markitdown", r.error.lower())


@unittest.skipUnless(_has_markitdown(), "markitdown not installed")
class TestMarkItDownIntegration(unittest.TestCase):
    """markitdown インストール環境でドキュメント変換が動作する。

    出力フォーマットは MarkItDown が決定するため、本文断片の存在のみ検証する
    （見出しの自動付与は撤去済み: 設計書 Q4）。
    """

    def test_html_via_markitdown(self) -> None:
        from hve.gui.doc_convert import convert_file

        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "page.html"
            src.write_text(
                "<html><body><h1>Hello World</h1><p>Body text.</p></body></html>",
                encoding="utf-8",
            )
            out_dir = Path(td) / "out"
            r = convert_file(src, out_dir=out_dir)
            self.assertTrue(r.ok, msg=f"error: {r.error}")
            assert r.converted_path is not None
            content = r.converted_path.read_text(encoding="utf-8")
            self.assertIn("Hello World", content)
            self.assertIn("Body text", content)

    def test_docx_via_markitdown(self) -> None:
        try:
            from docx import Document  # type: ignore[import-not-found]
        except ImportError:
            self.skipTest("python-docx not available (transitive of markitdown[all])")

        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "doc.docx"
            d = Document()
            d.add_heading("DocxHeading", level=1)
            d.add_paragraph("DocxBodyToken")
            d.save(str(src))

            from hve.gui.doc_convert import convert_file
            out_dir = Path(td) / "out"
            r = convert_file(src, out_dir=out_dir)
            self.assertTrue(r.ok, msg=f"error: {r.error}")
            assert r.converted_path is not None
            content = r.converted_path.read_text(encoding="utf-8")
            self.assertIn("DocxHeading", content)
            self.assertIn("DocxBodyToken", content)

    def test_xlsx_via_markitdown(self) -> None:
        try:
            from openpyxl import Workbook  # type: ignore[import-not-found]
        except ImportError:
            self.skipTest("openpyxl not available (transitive of markitdown[all])")

        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "book.xlsx"
            wb = Workbook()
            ws = wb.active
            ws.title = "Sheet1"
            ws.append(["ColA", "ColB"])
            ws.append(["XlsxTokenA", "XlsxTokenB"])
            wb.save(str(src))

            from hve.gui.doc_convert import convert_file
            out_dir = Path(td) / "out"
            r = convert_file(src, out_dir=out_dir)
            self.assertTrue(r.ok, msg=f"error: {r.error}")
            assert r.converted_path is not None
            content = r.converted_path.read_text(encoding="utf-8")
            self.assertIn("XlsxTokenA", content)
            self.assertIn("XlsxTokenB", content)

    def test_pptx_via_markitdown(self) -> None:
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
        except ImportError:
            self.skipTest("python-pptx not available (transitive of markitdown[all])")

        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "deck.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "PptxTitleToken"
            prs.save(str(src))

            from hve.gui.doc_convert import convert_file
            out_dir = Path(td) / "out"
            r = convert_file(src, out_dir=out_dir)
            self.assertTrue(r.ok, msg=f"error: {r.error}")
            assert r.converted_path is not None
            content = r.converted_path.read_text(encoding="utf-8")
            self.assertIn("PptxTitleToken", content)

    def test_xls_via_markitdown(self) -> None:
        try:
            import xlwt  # type: ignore[import-not-found]
        except ImportError:
            self.skipTest("xlwt not available (legacy .xls writer)")

        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "legacy.xls"
            wb = xlwt.Workbook()
            ws = wb.add_sheet("Sheet1")
            ws.write(0, 0, "XlsTokenA")
            ws.write(0, 1, "XlsTokenB")
            wb.save(str(src))

            from hve.gui.doc_convert import convert_file
            out_dir = Path(td) / "out"
            r = convert_file(src, out_dir=out_dir)
            self.assertTrue(r.ok, msg=f"error: {r.error}")
            assert r.converted_path is not None
            content = r.converted_path.read_text(encoding="utf-8")
            self.assertIn("XlsTokenA", content)


class TestChooseOriginFile(unittest.TestCase):
    """choose_origin_file は QApplication 不要な純関数。"""

    def test_empty_returns_none(self) -> None:
        try:
            from hve.gui.page_options_ard import choose_origin_file
        except ImportError:
            self.skipTest("PySide6 not installed")

        self.assertIsNone(choose_origin_file([]))

    def test_first_returned_when_multiple(self) -> None:
        try:
            from hve.gui.page_options_ard import choose_origin_file
        except ImportError:
            self.skipTest("PySide6 not installed")

        paths = [Path("a.md"), Path("b.md"), Path("c.md")]
        chosen = choose_origin_file(paths)
        self.assertEqual(chosen, paths[0])


if __name__ == "__main__":
    unittest.main()